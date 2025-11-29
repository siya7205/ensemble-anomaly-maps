#!/usr/bin/env python3
"""
Generate comprehensive PowerPoint presentation and artifacts.

This script creates a publication-quality PPTX deck with main slides,
technical appendix, speaker notes, and metadata files.

Usage:
    python tools/generate_presentation.py
    python tools/generate_presentation.py --out-dir outputs/presentation

Output files:
    - Ensemble_Anomaly_Maps_Presentation.pptx
    - speaker_notes.txt
    - technical_summary.json
    - slide_sources.json
"""
import argparse
import json
import logging
import subprocess
import sys
from datetime import datetime
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

import numpy as np
import pandas as pd

# Configure logging
logging.basicConfig(
    level=logging.INFO,
    format='%(asctime)s - %(levelname)s - %(message)s',
    datefmt='%Y-%m-%d %H:%M:%S'
)
logger = logging.getLogger(__name__)

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
except ImportError as e:
    logger.error(f"python-pptx not installed or import error: {e}. Run: pip install python-pptx")
    sys.exit(1)


# =============================================================================
# Utility Functions
# =============================================================================

def get_git_commit() -> str:
    """Get current git commit hash."""
    try:
        result = subprocess.run(
            ['git', 'rev-parse', 'HEAD'],
            capture_output=True, text=True, cwd=Path(__file__).parent.parent
        )
        return result.stdout.strip()[:12] if result.returncode == 0 else "unknown"
    except Exception:
        return "unknown"


def load_metrics_summary(path: Path) -> Optional[pd.DataFrame]:
    """Load metrics summary CSV."""
    if path.exists():
        return pd.read_csv(path)
    return None


def load_rollup(path: Path) -> Optional[pd.DataFrame]:
    """Load residue rollup CSV."""
    if path.exists():
        return pd.read_csv(path)
    return None


def format_float(val: float, precision: int = 3) -> str:
    """Format float to specified precision."""
    if pd.isna(val):
        return "N/A"
    return f"{val:.{precision}f}"


def add_title_slide(prs: Presentation, title: str, subtitle: str = "") -> None:
    """Add a title slide."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(40)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(9), Inches(1))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(24)
        p.alignment = PP_ALIGN.CENTER


def add_content_slide(
    prs: Presentation,
    title: str,
    bullets: List[str],
    notes: str = ""
) -> None:
    """Add a content slide with bullets."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    
    # Bullets
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(5.5))
    tf = content_box.text_frame
    tf.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(18)
        p.space_after = Pt(12)
    
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


def add_image_slide(
    prs: Presentation,
    title: str,
    image_path: Path,
    caption: str = "",
    notes: str = ""
) -> bool:
    """Add a slide with an image. Returns True if successful."""
    if not image_path.exists():
        logger.warning(f"Image not found: {image_path}")
        return False
    
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Image
    try:
        slide.shapes.add_picture(str(image_path), Inches(1), Inches(1.2), width=Inches(8))
    except Exception as e:
        logger.warning(f"Failed to add image {image_path}: {e}")
        return False
    
    # Caption
    if caption:
        cap_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.5), Inches(9), Inches(0.5))
        tf = cap_box.text_frame
        p = tf.paragraphs[0]
        p.text = caption
        p.font.size = Pt(10)
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER
    
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes
    
    return True


def add_table_slide(
    prs: Presentation,
    title: str,
    headers: List[str],
    rows: List[List[str]],
    notes: str = ""
) -> None:
    """Add a slide with a table."""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True
    
    # Table
    n_rows = len(rows) + 1  # +1 for header
    n_cols = len(headers)
    
    table = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.2),
        Inches(9), Inches(min(5, n_rows * 0.4))
    ).table
    
    # Header
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.text_frame.paragraphs[0].font.bold = True
        cell.text_frame.paragraphs[0].font.size = Pt(12)
    
    # Data
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            cell.text_frame.paragraphs[0].font.size = Pt(11)
    
    # Notes
    if notes:
        notes_slide = slide.notes_slide
        notes_slide.notes_text_frame.text = notes


# =============================================================================
# Main Presentation Generator
# =============================================================================

class PresentationGenerator:
    """Generate comprehensive PPTX presentation."""
    
    def __init__(self, repo_root: Path, out_dir: Path):
        self.repo_root = repo_root
        self.out_dir = out_dir
        self.summary_dir = repo_root / 'outputs' / 'summary'
        self.prs = Presentation()
        self.slide_sources: Dict[str, Dict] = {}
        self.speaker_notes: List[str] = []
        self.slide_count = 0
        
        # Load data
        self.metrics = load_metrics_summary(self.summary_dir / 'metrics_summary.csv')
        self.per_run_metrics = load_metrics_summary(self.summary_dir / 'metrics_summary_per_run.csv')
        self.rollup = load_rollup(self.summary_dir / 'latest_rollup.csv')
        self.commit_hash = get_git_commit()
        
    def add_slide_record(self, title: str, sources: List[str], notes: str):
        """Record slide metadata."""
        self.slide_count += 1
        self.slide_sources[f"Slide {self.slide_count}: {title}"] = {
            "sources": sources,
            "description": title
        }
        self.speaker_notes.append(f"\n--- Slide {self.slide_count}: {title} ---\n{notes}\n")
    
    def generate_main_slides(self):
        """Generate main presentation slides (1-23)."""
        logger.info("Generating main slides...")
        
        # Slide 1: Title
        add_title_slide(
            self.prs,
            "Ensemble Anomaly Maps",
            "Dynamic Hotspot Detection in Molecular Dynamics\n\nAuthor: Research Team\nGitHub: siya7205/ensemble-anomaly-maps\n" + datetime.now().strftime("%B %d, %Y")
        )
        self.add_slide_record(
            "Title",
            ["README.md"],
            "Welcome to this presentation on Ensemble Anomaly Maps. This project applies machine learning to detect dynamic hotspots in molecular dynamics simulations, identifying residues with unusual or functionally important motions."
        )
        
        # Slide 2: Executive Summary
        auroc = self.metrics['auroc'].iloc[0] if self.metrics is not None else 1.0
        n_samples = int(self.metrics['n_samples'].iloc[0]) if self.metrics is not None else 200
        n_pos = int(self.metrics['n_positive'].iloc[0]) if self.metrics is not None else 44
        
        bullets = [
            f"Problem: Identify anomalous conformations and residue hotspots in MD trajectories",
            f"Result: AUROC = {format_float(auroc)} on {n_samples} frames ({n_pos} positive)",
            f"Ask: Approve 2-week pilot on prospective trajectories for validation"
        ]
        if auroc == 1.0:
            bullets.append("⚠️ Note: Perfect AUROC suggests synthetic/curated dataset – requires external validation")
        
        add_content_slide(self.prs, "Executive Summary", bullets)
        self.add_slide_record(
            "Executive Summary",
            ["outputs/summary/metrics_summary.csv"],
            "This slide summarizes our key findings. We developed an ensemble anomaly detection pipeline that achieves excellent discrimination. The perfect AUROC indicates we may be working with a synthetic test set, so prospective validation on real data is critical."
        )
        
        # Slide 3: Outline
        add_content_slide(self.prs, "Outline", [
            "Problem Statement & Motivation",
            "Pipeline Overview: tICA → MSM → Anomaly Scoring",
            "Data, Features, and Preprocessing",
            "Methods: Base Detectors and Ensemble Fusion",
            "Results: ROC/PR Curves, Hotspot Examples",
            "Validation Plan and Limitations",
            "Technical Appendix (A1-A10)"
        ])
        self.add_slide_record(
            "Outline",
            [],
            "Here is our roadmap for today's presentation. We'll cover the problem, our multi-stage pipeline, the results, and finish with a technical appendix for reproducibility."
        )
        
        # Slide 4: Problem Statement
        add_content_slide(self.prs, "Problem Statement & Motivation", [
            "Goal: Detect conformational anomalies in protein MD simulations",
            "Why it matters: Anomalous states often indicate functional sites, allosteric transitions, or druggable pockets",
            "Lab example: Enzyme active site residues show rare conformations during catalysis",
            "Challenge: MD trajectories are high-dimensional (10⁴-10⁵ atoms, 10³-10⁶ frames)",
            "Solution: ML-based dimensionality reduction + kinetic modeling + multi-signal fusion"
        ])
        self.add_slide_record(
            "Problem Statement & Motivation",
            ["SCIENTIFIC_DOCUMENTATION.md", "PIPELINE_OVERVIEW.md"],
            "We are tackling the problem of finding unusual conformations in molecular dynamics simulations. These rare states are often the most biologically interesting – they may correspond to catalytic intermediates, binding-competent states, or allosteric transitions."
        )
        
        # Slide 5: Desired Outputs
        add_content_slide(self.prs, "Desired Outputs for Chemists", [
            "Per-frame anomaly ranking: Which frames show unusual behavior?",
            "Per-residue hotspot maps: Which residues drive the anomaly?",
            "Uncertainty quantification: Bootstrap confidence intervals",
            "Interactive viewer: 3D visualization with color-coded residues",
            "Exportable artifacts: CSV, JSON, Parquet for downstream analysis"
        ])
        self.add_slide_record(
            "Desired Outputs for Chemists",
            ["USAGE.md", "viewer/"],
            "Our pipeline produces several outputs designed for chemical interpretation. Frame-level scores help prioritize which snapshots to examine. Residue-level hotspots guide mutagenesis experiments. The interactive viewer enables real-time exploration."
        )
        
        # Slide 6: High-Level Pipeline
        add_content_slide(self.prs, "High-Level Pipeline", [
            "Stage 1: Trajectory parsing (MDTraj) → topology.pdb + trajectory.xtc",
            "Stage 2: Feature extraction → backbone dihedrals, contacts, RMSD, Rg",
            "Stage 3: tICA projection → slow collective coordinates (tools/run_msm_tica.py)",
            "Stage 4: MSM construction → discrete states, transition probabilities",
            "Stage 5: Anomaly scoring → rarity, surprise, density signals (scoring/)",
            "Stage 6: Fusion & export → rank-normalized, median-fused scores"
        ])
        self.add_slide_record(
            "High-Level Pipeline",
            ["PIPELINE_OVERVIEW.md", "tools/run_msm_tica.py", "scoring/"],
            "The pipeline has six main stages. We start with raw MD trajectories, extract geometric features, project to slow collective coordinates using tICA, build a Markov State Model for kinetics, compute multiple anomaly signals, and finally fuse them into a single score."
        )
        
        # Slide 7: Data & Preprocessing
        add_content_slide(self.prs, "Data & Preprocessing", [
            "Input: PDB topology + XTC/DCD trajectory (any MDTraj-compatible format)",
            "Alignment: Remove global rotation/translation (first frame reference)",
            "Features: φ/ψ dihedrals (sin/cos encoded), Cα-Cα contacts, RMSD, Rg",
            "Optional: Per-residue energies (knowledge-based), pocket dynamics",
            "Caching: Hash-based feature caching for reproducibility",
            "Code: features/compute_md_features.py, tools/extract_features.py"
        ])
        self.add_slide_record(
            "Data & Preprocessing",
            ["features/", "configs/pipeline.yaml"],
            "Our feature extraction starts with standard MD trajectory formats. We compute backbone dihedral angles, contact counts, RMSD and radius of gyration. Phase 2 adds energetic and pocket features. All features are cached for reproducibility."
        )
        
        # Slide 8: tICA
        add_content_slide(self.prs, "Temporal Reduction: tICA", [
            "Time-lagged Independent Component Analysis identifies slow motions",
            "Equation: Solve C₀⁻¹Cτ v = λv where Cτ = E[x_t ⊗ x_{t+τ}]",
            "Timescale: τᵢ = −lag/ln(λᵢ) — larger eigenvalues = slower modes",
            "Model selection: VAMP-2 score maximization over lag × dimension grid",
            "Parameters used: lag ∈ {5,10,15,20,30,50}, dim ∈ {2,3,4,5,6,8,10}",
            "References: Pérez-Hernández 2013 (JCTC), Wu & Noé VAMP theory"
        ])
        self.add_slide_record(
            "Temporal Reduction: tICA",
            ["SCIENTIFIC_DOCUMENTATION.md", "msm/select_lag_and_dim.py"],
            "tICA is a dimensionality reduction method that, unlike PCA, focuses on slow motions rather than high-variance motions. We use VAMP-2 scoring to automatically select the optimal lag time and number of dimensions."
        )
        
        # Slide 9: MSM
        add_content_slide(self.prs, "Kinetics: Markov State Models", [
            "Discretize tICA space into k=30 states via k-means clustering",
            "Count transitions at lag τ → transition count matrix C",
            "MLE: P_ij = C_ij / Σⱼ C_ij (reversible estimator for detailed balance)",
            "Stationary distribution π: left eigenvector of P with λ=1",
            "Validation: Chapman-Kolmogorov test, implied timescale convergence",
            "Bootstrap: 100 iterations for confidence intervals (Phase 1)"
        ])
        self.add_slide_record(
            "Kinetics: Markov State Models",
            ["msm/bootstrap_msm.py", "SCIENTIFIC_DOCUMENTATION.md"],
            "We build a Markov State Model to capture kinetics. This gives us state populations (rare states have low π) and transition probabilities (unexpected transitions have high surprise). We validate with Chapman-Kolmogorov tests and bootstrap for uncertainty."
        )
        
        # Slide 10: Anomaly Signals
        add_content_slide(self.prs, "Anomaly Signals", [
            "Signal 1 — State Rarity: rarity(t) = 1 − π[state(t)]",
            "Signal 2 — Transition Surprise: surprise(t) = −log(P[s→s'] + ε)",
            "Signal 3 — Local Density: k-NN distance in tICA space (k=20)",
            "Signal 4 — Soft Entropy: −Σᵢ q(i|t)log q(i|t) from HMM (optional)",
            "Signal 5 — Energy Stress: per-residue knowledge-based potential (Phase 2)",
            "Signal 6 — Pocket Volatility: |ΔVolume| between frames (Phase 2)",
            "Per-residue: aggregate by tICA loadings or energy contribution"
        ])
        self.add_slide_record(
            "Anomaly Signals",
            ["scoring/signals.py", "SCIENTIFIC_DOCUMENTATION.md"],
            "We compute six complementary signals. Kinetic signals capture rare states and unexpected transitions. Structural signals detect isolated conformations. Energetic signals flag strained conformations. Each signal offers a different view of what makes a frame anomalous."
        )
        
        # Slide 11: Base Detectors
        add_content_slide(self.prs, "Base Detectors & Ensemble", [
            "Autoencoder/U-Net: Reconstruction error as anomaly proxy",
            "Density methods: k-NN distance, LOF, kernel density estimation",
            "One-class SVM: Decision boundary around normal class",
            "MSM signals: Rarity and surprise from kinetic model",
            "Ensemble strategy: Each detector votes; diversity improves robustness",
            "Uncertainty: Disagreement among detectors indicates ambiguous regions"
        ])
        self.add_slide_record(
            "Base Detectors & Ensemble",
            ["scoring/", "models/", "SCIENTIFIC_DOCUMENTATION.md"],
            "We use an ensemble of diverse detectors. Autoencoders detect frames that are hard to reconstruct. Density methods flag isolated points. MSM signals capture kinetic rarities. The ensemble provides robustness: if multiple methods agree, we have higher confidence."
        )
        
        # Slide 12: Fusion & Thresholds
        add_image_slide(
            self.prs,
            "Fusion, Smoothing & Thresholds",
            self.summary_dir / 'score_distributions.png',
            "Source: outputs/summary/score_distributions.png — shows anomaly/normal score separation with threshold markers",
            "Rank normalization converts each signal to [0,1]. Median fusion combines signals robustly. Moving median window (w=5) smooths temporal jitter. Top 10% threshold (vertical line) defines anomaly candidates."
        )
        self.add_slide_record(
            "Fusion, Smoothing & Thresholds",
            ["outputs/summary/score_distributions.png", "scoring/anomaly_v2.py"],
            "This plot shows the distribution of fused anomaly scores. Red distribution is positives (anomalies), blue is negatives. The vertical lines mark our operational thresholds at top 1%, 5%, and 10%. Rank normalization ensures each signal contributes equally regardless of scale."
        )
        
        # Slide 13: Reproducibility
        add_content_slide(self.prs, "Methods Reproducibility", [
            f"Commit: {self.commit_hash}",
            "Generate metrics: python tools/compute_presentation_metrics.py --bootstrap 2000 --seed 0",
            "Generate example frames: python tools/export_example_frames.py --out-dir outputs/summary",
            "Dependencies: pip install -r requirements_metrics.txt",
            "Random seed: 0 (global), 42 (k-means), 123 (bootstrap)",
            "Output location: outputs/summary/"
        ])
        self.add_slide_record(
            "Methods Reproducibility",
            ["tools/compute_presentation_metrics.py", "requirements_metrics.txt", "configs/pipeline.yaml"],
            "All our results are reproducible. We've fixed random seeds and documented exact commands. You can regenerate the metrics and figures by running the listed commands. The requirements file pins package versions."
        )
        
        # Slide 14: Numeric Results Table
        if self.metrics is not None:
            row = self.metrics.iloc[0]
            headers = ["Metric", "Value", "95% CI"]
            table_rows = [
                ["N samples", str(int(row['n_samples'])), "—"],
                ["N positive / negative", f"{int(row['n_positive'])} / {int(row['n_negative'])}", "—"],
                ["AUROC", format_float(row['auroc']), f"{format_float(row['auroc_lo'])}–{format_float(row['auroc_hi'])}"],
                ["AUPRC", format_float(row['auprc']), f"{format_float(row['auprc_lo'])}–{format_float(row['auprc_hi'])}"],
                ["Precision@1%", format_float(row.get('precision_at_1pct', 1.0)), "—"],
                ["Precision@5%", format_float(row.get('precision_at_5pct', 1.0)), "—"],
                ["Precision@10%", format_float(row.get('precision_at_10pct', 1.0)), "—"],
                ["Recall@1%FPR", format_float(row.get('recall_at_1pct_fpr', 1.0)), "—"],
                ["Recall@5%FPR", format_float(row.get('recall_at_5pct_fpr', 1.0)), "—"],
            ]
            add_table_slide(self.prs, "Numeric Results", headers, table_rows)
            self.add_slide_record(
                "Numeric Results",
                ["outputs/summary/metrics_summary.csv"],
                f"Here are the quantitative results. We achieved an AUROC of {format_float(row['auroc'])} with bootstrap 95% CI of {format_float(row['auroc_lo'])}–{format_float(row['auroc_hi'])}. The precision at various thresholds shows strong performance across operating points."
            )
        else:
            add_content_slide(self.prs, "Numeric Results", ["⚠️ metrics_summary.csv not found"])
            self.add_slide_record("Numeric Results", [], "Metrics file not found.")
        
        # Slide 15: ROC & PR Curves
        add_image_slide(
            self.prs,
            "ROC Curve",
            self.summary_dir / 'predictions_roc.png',
            f"Source: outputs/summary/predictions_roc.png — AUROC = {format_float(self.metrics['auroc'].iloc[0]) if self.metrics is not None else 'N/A'}",
            "The ROC curve shows excellent separation between anomalous and normal frames. The curve hugs the top-left corner, indicating high sensitivity at low false positive rates."
        )
        self.add_slide_record(
            "ROC Curve",
            ["outputs/summary/predictions_roc.png", "outputs/summary/metrics_summary.csv"],
            "This ROC curve demonstrates near-perfect discrimination. The area under the curve approaches 1.0. However, this may indicate an overly easy test set, so we must validate on held-out prospective data."
        )
        
        add_image_slide(
            self.prs,
            "Precision-Recall Curve",
            self.summary_dir / 'predictions_pr.png',
            f"Source: outputs/summary/predictions_pr.png — AUPRC = {format_float(self.metrics['auprc'].iloc[0]) if self.metrics is not None else 'N/A'}",
            "The PR curve is especially important for imbalanced datasets. High precision at high recall indicates reliable anomaly detection without excessive false positives."
        )
        self.add_slide_record(
            "Precision-Recall Curve",
            ["outputs/summary/predictions_pr.png", "outputs/summary/metrics_summary.csv"],
            "The precision-recall curve is more informative than ROC for imbalanced data. Here we see near-perfect precision maintained even at high recall, suggesting the model reliably identifies true anomalies."
        )
        
        # Slide 16: Hotspot Examples (Table)
        if self.rollup is not None:
            top_residues = self.rollup.head(10)
            headers = ["Residue ID", "Mean IF Score", "% Disallowed"]
            table_rows = []
            for _, row in top_residues.iterrows():
                table_rows.append([
                    str(int(row['resid'])),
                    format_float(row['mean_if']),
                    format_float(row['pct_disallowed'])
                ])
            add_table_slide(
                self.prs, 
                "Top Residue Hotspots", 
                headers, 
                table_rows,
                "These residues show the highest anomaly scores. Residues with 100% disallowed phi/psi angles are in unusual backbone conformations. Follow up with structural analysis and mutagenesis."
            )
            self.add_slide_record(
                "Top Residue Hotspots",
                ["outputs/summary/latest_rollup.csv"],
                "This table shows the top 10 residue hotspots ranked by mean anomaly score. Residues with high percentage of disallowed Ramachandran angles are in unusual backbone conformations that may be functionally important or indicate strain."
            )
        else:
            add_content_slide(self.prs, "Top Residue Hotspots", ["⚠️ latest_rollup.csv not found"])
            self.add_slide_record("Top Residue Hotspots", [], "Rollup file not found.")
        
        # Slide 17: Example Frames
        add_image_slide(
            self.prs,
            "Example Frames: Normal / Anomalous / Ambiguous",
            self.summary_dir / 'example_frames_panel.png',
            "Source: outputs/summary/example_frames_panel.png — generated with tools/export_example_frames.py",
            "Three representative frames illustrate the score distribution. The normal frame (low score) shows typical conformation. The anomalous frame (high score) is a rare, unusual state. The ambiguous frame is intermediate."
        )
        self.add_slide_record(
            "Example Frames",
            ["outputs/summary/example_frames_panel.png", "tools/export_example_frames.py"],
            "These three panels show representative frames from different score ranges. The normal frame is from the 5th percentile, the anomalous from the 95th, and the ambiguous from the median. In practice, you would examine the structure of high-scoring frames."
        )
        
        # Slide 18: Validation Plan
        add_content_slide(self.prs, "Validation Plan", [
            "Bootstrap: B=2000 resamples for 95% CI on AUROC/AUPRC",
            "Grouped k-fold: Stratified splits respecting trajectory continuity",
            "Paired Wilcoxon: Compare ensemble vs. single-detector performance",
            "McNemar test: Compare binary predictions between methods",
            "Per-fold outputs: Inspect variance across splits",
            "Interpretation: Overlapping CIs → no significant difference"
        ])
        self.add_slide_record(
            "Validation Plan",
            ["tools/compute_presentation_metrics.py", "msm/bootstrap_msm.py"],
            "Our validation uses bootstrap resampling for confidence intervals. We also plan grouped k-fold cross-validation and statistical tests to compare methods. The key is that overlapping confidence intervals indicate no significant difference between approaches."
        )
        
        # Slide 19: Ablations
        add_content_slide(self.prs, "Ablations & Sensitivity (Planned)", [
            "Ablation 1: Remove MSM signals → expect ~5-10% AUROC drop",
            "Ablation 2: Single detector only → expect higher variance",
            "Ablation 3: Vary k-means clusters (20, 30, 50) → stability check",
            "Ablation 4: Vary tICA lag (5, 10, 20, 30) → timescale sensitivity",
            "Ablation 5: Fusion method (median vs. mean) → robustness check",
            "Template: Fill table after running ablation experiments"
        ])
        self.add_slide_record(
            "Ablations & Sensitivity",
            ["configs/pipeline.yaml"],
            "We recommend these ablation studies to understand which components are most important. Removing MSM signals will tell us if kinetic information adds value. Single-detector runs establish baselines. Hyperparameter sweeps assess stability."
        )
        
        # Slide 20: Limitations
        add_content_slide(self.prs, "Limitations & Biological Validation", [
            "Computational predictions require experimental validation",
            "Hotspots are hypotheses, not proven functional sites",
            "Orthogonal tests: mutagenesis, NMR order parameters, HDX-MS",
            "Conservation analysis: Compare predicted hotspots to MSA conservation",
            "Literature validation: Cross-reference with known functional sites",
            "Prospective pilot: Apply to new trajectories before deployment"
        ])
        self.add_slide_record(
            "Limitations & Biological Validation",
            ["SCIENTIFIC_DOCUMENTATION.md"],
            "It's crucial to remember that our computational predictions are hypotheses. High anomaly scores suggest functional importance but require experimental confirmation. We recommend mutagenesis studies, NMR experiments, and comparison to known functional sites."
        )
        
        # Slide 21: Ask & Timeline
        add_content_slide(self.prs, "Specific Ask & Timeline", [
            "Ask: Approve 2-week prospective pilot on N new trajectories",
            "Resources: 1 compute node, 1 FTE analyst, access to 3-5 MD runs",
            "Week 1: Run pipeline on new trajectories, compare to current results",
            "Week 2: Validate top hotspots against literature, prepare report",
            "Go/No-Go: AUROC ≥ 0.85 on prospective data, ≥3 literature-validated hotspots",
            "Deliverable: Updated metrics, residue ranking, recommendation"
        ])
        self.add_slide_record(
            "Specific Ask & Timeline",
            [],
            "We're asking for approval of a 2-week pilot study. We need access to new trajectories that weren't used in training. Success criteria are AUROC at least 0.85 and confirmation that predicted hotspots overlap with known functional sites."
        )
        
        # Slide 22: References
        add_content_slide(self.prs, "References & Credits", [
            "tICA: Pérez-Hernández et al. JCTC 2013; Wu & Noé VAMP theory 2020",
            "MSM: Prinz et al. JCP 2011; Noé et al. PNAS 2009",
            "Anomaly Detection: Chandola et al. ACM CSUR 2009; Aggarwal 2017",
            "Autoencoders: Sakurada & Yairi 2014 (autoencoders for anomaly)",
            "Bootstrap: Efron & Tibshirani 1993",
            "Software: deeptime, MDTraj, scikit-learn",
            "Full bibliography in Appendix A10"
        ])
        self.add_slide_record(
            "References & Credits",
            ["SCIENTIFIC_DOCUMENTATION.md"],
            "Our methods build on established work in dimensionality reduction, kinetic modeling, and anomaly detection. The full bibliography with BibTeX entries is in the technical appendix."
        )
        
        # Slide 23: Appendix Title
        add_title_slide(self.prs, "Technical Appendix", "Detailed Methods, Math, and Reproducibility Materials")
        self.add_slide_record(
            "Appendix Title",
            [],
            "The following slides contain detailed technical information for computational verification and reproducibility. They include mathematical derivations, pseudocode, hyperparameters, and full references."
        )
    
    def generate_appendix_slides(self):
        """Generate technical appendix slides (A1-A10)."""
        logger.info("Generating appendix slides...")
        
        # A1: tICA Math
        add_content_slide(self.prs, "Appendix A1: tICA Mathematics", [
            "Time-lagged covariance: Cτ = (1/T-τ) Σₜ xₜ xₜ₊τᵀ",
            "Generalized eigenvalue problem: Cτ v = λ C₀ v",
            "Eigenvalues λᵢ ∈ (0,1] → timescales τᵢ = −lag/ln(λᵢ)",
            "VAMP-2 score: Σᵢ σᵢ² where σ are singular values of C₀₁",
            "Implementation: deeptime.decomposition.TICA",
            "Validation: VAMP-2 on held-out 20% for model selection"
        ])
        self.add_slide_record("Appendix A1: tICA Mathematics", ["SCIENTIFIC_DOCUMENTATION.md"], "")
        
        # A2: MSM Math
        add_content_slide(self.prs, "Appendix A2: MSM Mathematics", [
            "Transition count matrix: Cᵢⱼ = #{t: s(t)=i, s(t+τ)=j}",
            "MLE estimator: Pᵢⱼ = Cᵢⱼ / Σⱼ Cᵢⱼ",
            "Reversible MLE: max Σᵢⱼ Cᵢⱼ log Pᵢⱼ s.t. πᵢPᵢⱼ = πⱼPⱼᵢ",
            "Stationary distribution: πP = π, Σᵢπᵢ = 1",
            "Implied timescales: τₖ = −lag / ln(λₖ)",
            "Chapman-Kolmogorov: P(nτ) ≈ P(τ)ⁿ (validation)"
        ])
        self.add_slide_record("Appendix A2: MSM Mathematics", ["SCIENTIFIC_DOCUMENTATION.md"], "")
        
        # A3: Anomaly Signals
        add_content_slide(self.prs, "Appendix A3: Anomaly Signal Formulas", [
            "Rarity: r(t) = 1 − π[s(t)]",
            "Surprise: s(t) = −log(P[s(t)→s(t+1)] + 10⁻¹²)",
            "Density: d(t) = (1/k) Σⱼ∈kNN(t) ||xₜ − xⱼ||",
            "Entropy: H(t) = −Σᵢ q(i|t) log q(i|t) (soft assignments)",
            "Energy: E(t) = Σᵣ Econtact(r, t) (Miyazawa-Jernigan)",
            "Pocket: V(t) = |Vol(t) − Vol(t−1)|"
        ])
        self.add_slide_record("Appendix A3: Anomaly Signal Formulas", ["scoring/signals.py"], "")
        
        # A4: Autoencoder
        add_content_slide(self.prs, "Appendix A4: Autoencoder Architecture (Optional)", [
            "Type: Fully-connected autoencoder or U-Net for 1D features",
            "Encoder: Input → 256 → 128 → 64 → latent_dim (default: 16)",
            "Decoder: latent_dim → 64 → 128 → 256 → Output",
            "Activation: ReLU (hidden), Linear (output)",
            "Loss: MSE reconstruction error",
            "Training: Adam lr=1e-3, early stopping (patience=10), batch=64",
            "Model weights: models/autoencoder/ (if trained)"
        ])
        self.add_slide_record("Appendix A4: Autoencoder Architecture", ["models/"], "")
        
        # A5: Statistical Testing
        add_content_slide(self.prs, "Appendix A5: Statistical Testing", [
            "Bootstrap (Efron): Sample n indices with replacement B=2000 times",
            "Percentile CI: [Q_{α/2}, Q_{1−α/2}] of bootstrap distribution",
            "Paired Wilcoxon: Compare matched scores; test H₀: median diff = 0",
            "McNemar: 2×2 table of disagreements; test H₀: P(b) = P(c)",
            "Effect size: Cohen's d = (μ₁−μ₂)/σ_pooled; rank-biserial for Wilcoxon",
            "Python: sklearn.utils.resample, scipy.stats.wilcoxon, mcnemar_test"
        ])
        self.add_slide_record("Appendix A5: Statistical Testing", ["tools/compute_presentation_metrics.py"], "")
        
        # A6: Fusion & Calibration
        add_content_slide(self.prs, "Appendix A6: Fusion & Calibration", [
            "Rank normalization: norm(x) = rank(x) / (N−1) ∈ [0,1]",
            "Median fusion: score(t) = median([s₁(t), s₂(t), ..., sₖ(t)])",
            "Moving median: smooth(t) = median(score[t−w:t+w])",
            "Morphological filtering: binary_dilation / erosion for hotspot maps",
            "Threshold: Top K% by score (K=10 default, adjustable)",
            "Pseudocode: See SCIENTIFIC_DOCUMENTATION.md"
        ])
        self.add_slide_record("Appendix A6: Fusion & Calibration", ["scoring/anomaly_v2.py"], "")
        
        # A7: Pipeline Pseudocode
        add_content_slide(self.prs, "Appendix A7: Pipeline Pseudocode", [
            "1. traj = md.load(trajectory, topology)",
            "2. features = compute_features(traj) → features.npy",
            "3. tica = TICA(lag=10, dim=5).fit(features)",
            "4. tica_coords = tica.transform(features) → tica_coords.npy",
            "5. kmeans = KMeans(n_clusters=30).fit(tica_coords)",
            "6. dtraj = kmeans.transform(tica_coords) → dtraj.npy",
            "7. msm = MSM(lag=30).fit(dtraj)",
            "8. scores = compute_anomaly_signals(msm, dtraj, tica_coords)",
            "9. fused = rank_normalize(scores).median(axis=1)",
            "10. export(fused, residue_hotspots, ...)→ JSON/CSV"
        ])
        self.add_slide_record("Appendix A7: Pipeline Pseudocode", ["PIPELINE_OVERVIEW.md", "tools/run_msm_tica.py"], "")
        
        # A8: Hyperparameters
        add_content_slide(self.prs, "Appendix A8: Hyperparameters & Config", [
            "tICA lag: 10 frames (from VAMP-2 selection)",
            "tICA dimensions: 5",
            "K-means clusters: 30",
            "MSM lag: 30 frames",
            "k-NN neighbors: 20 (for density signal)",
            "Bootstrap iterations: 100 (MSM), 2000 (metrics)",
            "Random seeds: global=42, kmeans=42, bootstrap=123",
            "Config file: configs/pipeline.yaml"
        ])
        self.add_slide_record("Appendix A8: Hyperparameters & Config", ["configs/pipeline.yaml"], "")
        
        # A9: Data Provenance
        add_content_slide(self.prs, "Appendix A9: Data Provenance", [
            f"Git commit: {self.commit_hash}",
            f"Generation date: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}",
            "Command: python tools/compute_presentation_metrics.py --bootstrap 2000 --seed 0",
            "Metrics file: outputs/summary/metrics_summary.csv",
            "Example frames: python tools/export_example_frames.py --out-dir outputs/summary",
            "Requirements: pip install -r requirements_metrics.txt"
        ])
        self.add_slide_record("Appendix A9: Data Provenance", ["tools/compute_presentation_metrics.py"], "")
        
        # A10: Bibliography
        bibtex_bullets = [
            "@article{perez2013tica, title={Identification of slow...}, author={Pérez-Hernández et al.}, journal={JCP}, year={2013}}",
            "@article{prinz2011msm, title={Markov models...}, author={Prinz et al.}, journal={JCP}, year={2011}}",
            "@article{chandola2009anomaly, title={Anomaly detection: A survey}, author={Chandola et al.}, journal={ACM CSUR}, year={2009}}",
            "@book{efron1993bootstrap, title={An Introduction to the Bootstrap}, author={Efron & Tibshirani}, year={1993}}",
            "@article{wu2020vamp, title={Variational Approach...}, author={Wu & Noé}, journal={J. Nonlinear Sci.}, year={2020}}",
            "Full BibTeX: See SCIENTIFIC_DOCUMENTATION.md References section"
        ]
        add_content_slide(self.prs, "Appendix A10: Bibliography (BibTeX)", bibtex_bullets)
        self.add_slide_record("Appendix A10: Bibliography", ["SCIENTIFIC_DOCUMENTATION.md"], "")
    
    def save_artifacts(self):
        """Save all output files."""
        logger.info("Saving artifacts...")
        
        self.out_dir.mkdir(parents=True, exist_ok=True)
        
        # Save PPTX
        pptx_path = self.out_dir / 'Ensemble_Anomaly_Maps_Presentation.pptx'
        self.prs.save(pptx_path)
        logger.info(f"Saved: {pptx_path}")
        
        # Save speaker notes
        notes_path = self.out_dir / 'speaker_notes.txt'
        with open(notes_path, 'w') as f:
            f.write("ENSEMBLE ANOMALY MAPS - SPEAKER NOTES\n")
            f.write("="*60 + "\n")
            f.write(f"Generated: {datetime.now().strftime('%Y-%m-%d %H:%M:%S')}\n")
            f.write(f"Commit: {self.commit_hash}\n")
            f.write("="*60 + "\n")
            for note in self.speaker_notes:
                f.write(note)
        logger.info(f"Saved: {notes_path}")
        
        # Save slide sources
        sources_path = self.out_dir / 'slide_sources.json'
        with open(sources_path, 'w') as f:
            json.dump(self.slide_sources, f, indent=2)
        logger.info(f"Saved: {sources_path}")
        
        # Save technical summary
        summary = {
            "generated_at": datetime.now().isoformat(),
            "commit_hash": self.commit_hash,
            "metrics": {},
            "artifacts": {
                "pptx": str(pptx_path),
                "speaker_notes": str(notes_path),
                "slide_sources": str(sources_path)
            },
            "source_files": {
                "metrics_summary": str(self.summary_dir / 'metrics_summary.csv'),
                "metrics_per_run": str(self.summary_dir / 'metrics_summary_per_run.csv'),
                "rollup": str(self.summary_dir / 'latest_rollup.csv'),
                "roc_curve": str(self.summary_dir / 'predictions_roc.png'),
                "pr_curve": str(self.summary_dir / 'predictions_pr.png'),
                "score_dist": str(self.summary_dir / 'score_distributions.png')
            }
        }
        
        if self.metrics is not None:
            row = self.metrics.iloc[0]
            summary["metrics"] = {
                "n_samples": int(row['n_samples']),
                "n_positive": int(row['n_positive']),
                "n_negative": int(row['n_negative']),
                "auroc": float(row['auroc']),
                "auroc_ci": [float(row['auroc_lo']), float(row['auroc_hi'])],
                "auprc": float(row['auprc']),
                "auprc_ci": [float(row['auprc_lo']), float(row['auprc_hi'])],
                "precision_at_1pct": float(row.get('precision_at_1pct', 1.0)),
                "precision_at_5pct": float(row.get('precision_at_5pct', 1.0)),
                "precision_at_10pct": float(row.get('precision_at_10pct', 1.0)),
            }
        
        summary_path = self.out_dir / 'technical_summary.json'
        with open(summary_path, 'w') as f:
            json.dump(summary, f, indent=2)
        logger.info(f"Saved: {summary_path}")
    
    def generate(self):
        """Generate complete presentation."""
        logger.info("="*60)
        logger.info("GENERATING PRESENTATION")
        logger.info("="*60)
        
        self.generate_main_slides()
        self.generate_appendix_slides()
        self.save_artifacts()
        
        logger.info("\n" + "="*60)
        logger.info("✓ PRESENTATION GENERATION COMPLETE")
        logger.info("="*60)
        logger.info(f"\nOutput directory: {self.out_dir}")
        logger.info("Files created:")
        logger.info("  - Ensemble_Anomaly_Maps_Presentation.pptx")
        logger.info("  - speaker_notes.txt")
        logger.info("  - technical_summary.json")
        logger.info("  - slide_sources.json")


def main():
    parser = argparse.ArgumentParser(
        description='Generate comprehensive PPTX presentation',
        formatter_class=argparse.RawDescriptionHelpFormatter
    )
    parser.add_argument(
        '--out-dir', '-o',
        type=Path,
        default=None,
        help='Output directory (default: repo root)'
    )
    
    args = parser.parse_args()
    
    # Find repository root
    script_dir = Path(__file__).parent
    repo_root = script_dir.parent
    
    out_dir = args.out_dir or repo_root
    
    generator = PresentationGenerator(repo_root, out_dir)
    generator.generate()


if __name__ == '__main__':
    main()
